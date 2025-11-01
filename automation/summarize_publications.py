"""Automation script for summarizing Tiede ja Teknologia resources.

This script fetches publications from the avoindata.fi CKAN API, extracts text
from PDF and PowerPoint resources, asks an OpenAI model to summarise the
content, and exports the summary as a PDF file. Generated summaries and the
processed resource list are stored locally so that the automation can skip
previously processed publications.

The script is intentionally written to be friendly for GitHub Actions use: the
state file lives inside the repository, the output folder can be versioned, and
runtime configuration happens via environment variables and command-line
arguments.
"""

from __future__ import annotations

import argparse
import json
import logging
import os
import re
import sys
import textwrap
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Dict, Iterable, Iterator, List, Optional

import requests
from fpdf import FPDF
from pptx import Presentation
from pptx.exc import PackageNotFoundError
from pypdf import PdfReader
from tenacity import retry, retry_if_exception_type, stop_after_attempt, wait_exponential

CKAN_API_BASE = "https://www.avoindata.fi/data/api/3/action"
SUPPORTED_FORMATS = {"pdf", "ppt", "pptx"}
DEFAULT_MODEL = os.getenv("OPENAI_SUMMARY_MODEL", "gpt-4o-mini")
OPENAI_API_URL = "https://api.openai.com/v1/chat/completions"


class AutomationError(RuntimeError):
    """Generic runtime error for the automation."""


@dataclass
class PublicationResource:
    """Represents a downloadable publication resource from CKAN."""

    dataset_id: str
    dataset_title: str
    dataset_notes: str | None
    resource_id: str
    resource_name: str
    resource_format: str
    resource_url: str
    last_modified: str | None

    @property
    def extension(self) -> str:
        fmt = self.resource_format.lower()
        if fmt == "ppt":
            return "ppt"
        if fmt == "pptx":
            return "pptx"
        return "pdf"


def setup_logging(verbose: bool) -> None:
    """Configure logging output."""

    level = logging.DEBUG if verbose else logging.INFO
    logging.basicConfig(
        level=level,
        format="%(asctime)s [%(levelname)s] %(message)s",
        datefmt="%Y-%m-%d %H:%M:%S",
    )


def slugify(value: str) -> str:
    """Create a filesystem-friendly slug."""

    value = value.lower()
    value = re.sub(r"[^a-z0-9]+", "-", value)
    value = value.strip("-")
    return value or "publication"


def load_state(path: Path) -> Dict[str, Dict[str, str]]:
    """Load the processed resource state file."""

    if not path.exists():
        logging.debug("State file %s does not exist yet", path)
        return {}

    try:
        with path.open("r", encoding="utf-8") as fh:
            data = json.load(fh)
    except json.JSONDecodeError as exc:
        raise AutomationError(f"State file {path} is corrupted: {exc}") from exc
    if not isinstance(data, dict):
        raise AutomationError(f"State file {path} is expected to contain a JSON object")
    return data


def save_state(path: Path, state: Dict[str, Dict[str, str]]) -> None:
    """Persist the processed resource state to disk."""

    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("w", encoding="utf-8") as fh:
        json.dump(state, fh, indent=2, ensure_ascii=False, sort_keys=True)
        fh.write("\n")


def fetch_group_packages(group_id: str) -> List[Dict[str, object]]:
    """Fetch all packages for the given CKAN group id."""

    url = f"{CKAN_API_BASE}/group_show"
    params = {"id": group_id, "include_datasets": True}
    logging.debug("Fetching packages for group %s", group_id)
    response = requests.get(url, params=params, timeout=60)
    response.raise_for_status()
    payload = response.json()
    if not payload.get("success"):
        raise AutomationError(f"CKAN API did not return success for group {group_id!r}")
    result = payload.get("result")
    if not isinstance(result, dict):
        raise AutomationError("Unexpected CKAN API payload shape")
    packages = result.get("packages")
    if not isinstance(packages, list):
        raise AutomationError("CKAN API returned unexpected package list")
    logging.info("Fetched %s packages for group %s", len(packages), group_id)
    return packages


def iter_publication_resources(packages: Iterable[Dict[str, object]]) -> Iterator[PublicationResource]:
    """Yield publication resources from CKAN packages."""

    for package in packages:
        if not isinstance(package, dict):
            continue
        dataset_id = str(package.get("id"))
        dataset_title = str(package.get("title")) if package.get("title") else dataset_id
        dataset_notes = package.get("notes") if isinstance(package.get("notes"), str) else None
        resources = package.get("resources")
        if not isinstance(resources, list):
            continue
        for resource in resources:
            if not isinstance(resource, dict):
                continue
            resource_format = str(resource.get("format", "")).lower()
            if resource_format not in SUPPORTED_FORMATS:
                continue
            resource_id = str(resource.get("id"))
            resource_name = str(resource.get("name") or resource.get("title") or dataset_title)
            resource_url = str(resource.get("url") or resource.get("download_url") or "")
            if not resource_url:
                logging.debug("Skipping resource %s without URL", resource_id)
                continue
            last_modified = resource.get("last_modified") if isinstance(resource.get("last_modified"), str) else None
            yield PublicationResource(
                dataset_id=dataset_id,
                dataset_title=dataset_title,
                dataset_notes=dataset_notes,
                resource_id=resource_id,
                resource_name=resource_name,
                resource_format=resource_format,
                resource_url=resource_url,
                last_modified=last_modified,
            )


def download_resource(resource: PublicationResource, destination_dir: Path) -> Path:
    """Download the resource file to the destination directory."""

    destination_dir.mkdir(parents=True, exist_ok=True)
    filename = f"{slugify(resource.resource_name)[:80]}-{resource.resource_id}.{resource.extension}"
    path = destination_dir / filename
    logging.info("Downloading %s", resource.resource_url)
    with requests.get(resource.resource_url, stream=True, timeout=120) as response:
        response.raise_for_status()
        with path.open("wb") as fh:
            for chunk in response.iter_content(chunk_size=8192):
                if chunk:
                    fh.write(chunk)
    return path


def extract_text_from_pdf(path: Path) -> str:
    """Extract text from a PDF file."""

    reader = PdfReader(str(path))
    text_parts: List[str] = []
    for page in reader.pages:
        extracted = page.extract_text() or ""
        if extracted:
            text_parts.append(extracted)
    return "\n".join(text_parts)


def extract_text_from_ppt(path: Path) -> str:
    """Extract text from a PowerPoint file."""

    try:
        presentation = Presentation(str(path))
    except PackageNotFoundError as exc:
        raise AutomationError(f"Could not open PowerPoint file {path}") from exc

    text_runs: List[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_runs.append(shape.text)
    return "\n".join(text_runs)


def extract_text(resource_path: Path, resource: PublicationResource) -> str:
    """Extract text content based on the resource format."""

    if resource.extension == "pdf":
        return extract_text_from_pdf(resource_path)
    return extract_text_from_ppt(resource_path)


def truncate_text(text: str, limit: int = 12000) -> str:
    """Truncate text so that we do not exceed the token budget."""

    if len(text) <= limit:
        return text
    logging.debug("Truncating extracted text from %s characters to %s", len(text), limit)
    return text[:limit]


def build_prompt(resource: PublicationResource, extracted_text: str) -> str:
    """Build the prompt for the OpenAI model."""

    meta_lines = [
        f"Dataset title: {resource.dataset_title}",
        f"Resource name: {resource.resource_name}",
        f"Resource format: {resource.resource_format.upper()}",
    ]
    if resource.last_modified:
        meta_lines.append(f"Last modified: {resource.last_modified}")
    if resource.dataset_notes:
        meta_lines.append("Dataset description: " + resource.dataset_notes.strip())

    prompt = textwrap.dedent(
        """
        You are an assistant that creates concise Finnish-language summaries of scientific
        and technology related publications. Read the provided content and write a
        structured summary that includes the following sections:
        - Tiivistelmä: 3-4 virkkeinen yleiskuva.
        - Keskeiset havainnot: luettelona 3-6 tärkeintä asiaa.
        - Julkaisun tiedot: yksityiskohtia kuten julkaisuvuosi, tekijät tai organisaatio
          jos ne käyvät ilmi lähdetekstistä.

        Jos jokin tieto ei ole saatavilla, mainitse siitä lyhyesti mutta älä keksi
        faktoja. Kirjoita selkeää asiantuntijasuomea.
        """
    ).strip()

    combined = "\n".join(meta_lines) + "\n\nLähteen sisältö alkaa:\n" + extracted_text
    trimmed = truncate_text(combined)
    return prompt + "\n\n" + trimmed


class OpenAIClient:
    """Simple OpenAI client for chat completions."""

    def __init__(self, api_key: str, model: str = DEFAULT_MODEL) -> None:
        if not api_key:
            raise AutomationError("OPENAI_API_KEY environment variable is not set")
        self.api_key = api_key
        self.model = model

    @retry(
        retry=retry_if_exception_type((requests.RequestException, AutomationError)),
        wait=wait_exponential(multiplier=2, min=2, max=30),
        stop=stop_after_attempt(5),
        reraise=True,
    )
    def create_summary(self, prompt: str) -> str:
        headers = {
            "Authorization": f"Bearer {self.api_key}",
            "Content-Type": "application/json",
        }
        payload = {
            "model": self.model,
            "messages": [
                {"role": "system", "content": "You create Finnish summaries of documents."},
                {"role": "user", "content": prompt},
            ],
            "temperature": 0.2,
        }
        logging.debug("Requesting summary from OpenAI model %s", self.model)
        response = requests.post(OPENAI_API_URL, headers=headers, json=payload, timeout=120)
        response.raise_for_status()
        data = response.json()
        choices = data.get("choices")
        if not choices:
            raise AutomationError("OpenAI API returned no choices")
        message = choices[0].get("message", {})
        content = message.get("content")
        if not content:
            raise AutomationError("OpenAI API response did not contain summary content")
        return str(content).strip()


def write_summary_pdf(summary_text: str, output_path: Path, metadata: PublicationResource) -> None:
    """Write the summary text to a PDF file."""

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Helvetica", "", 16)
    pdf.multi_cell(0, 10, metadata.dataset_title)
    pdf.ln(4)
    pdf.set_font("Helvetica", "", 12)
    for line in summary_text.splitlines():
        if not line.strip():
            pdf.ln(4)
            continue
        pdf.multi_cell(0, 6, line)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    pdf.output(str(output_path))


def process_resource(
    resource: PublicationResource,
    downloads_dir: Path,
    output_dir: Path,
    client: OpenAIClient,
) -> str:
    """Download, summarise, and persist the summary for a resource."""

    download_path = download_resource(resource, downloads_dir)
    try:
        extracted_text = extract_text(download_path, resource)
    finally:
        try:
            download_path.unlink()
        except OSError:
            logging.debug("Failed to delete temporary file %s", download_path)

    if not extracted_text.strip():
        logging.warning("Resource %s produced no text, using metadata only", resource.resource_id)
        extracted_text = resource.dataset_notes or resource.dataset_title

    prompt = build_prompt(resource, extracted_text)
    summary_text = client.create_summary(prompt)

    output_filename = f"{slugify(resource.dataset_title)[:60]}-{resource.resource_id}.pdf"
    output_path = output_dir / output_filename
    write_summary_pdf(summary_text, output_path, resource)
    logging.info("Written summary to %s", output_path)
    return output_filename


def parse_args(argv: Optional[List[str]] = None) -> argparse.Namespace:
    """Parse command line arguments."""

    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--group-id", default="tiede-ja-teknologia", help="CKAN group id to process")
    parser.add_argument("--output-dir", default="summaries", help="Where to store summary PDFs")
    parser.add_argument(
        "--state-file",
        default="data/processed_resources.json",
        help="Path to the JSON state file",
    )
    parser.add_argument(
        "--downloads-dir",
        default="downloads",
        help="Temporary directory for resource downloads",
    )
    parser.add_argument("--limit", type=int, default=None, help="Optional maximum number of new resources to process")
    parser.add_argument("--verbose", action="store_true", help="Enable debug logging")
    return parser.parse_args(argv)


def main(argv: Optional[List[str]] = None) -> int:
    args = parse_args(argv)
    setup_logging(args.verbose)

    state_path = Path(args.state_file)
    output_dir = Path(args.output_dir)
    downloads_dir = Path(args.downloads_dir)

    state = load_state(state_path)
    client = OpenAIClient(api_key=os.getenv("OPENAI_API_KEY"))

    packages = fetch_group_packages(args.group_id)
    processed_count = 0
    state_changed = False

    for resource in iter_publication_resources(packages):
        if resource.resource_id in state:
            logging.debug("Skipping already processed resource %s", resource.resource_id)
            continue
        try:
            summary_filename = process_resource(resource, downloads_dir, output_dir, client)
        except Exception as exc:  # pylint: disable=broad-except
            logging.exception("Failed to process resource %s: %s", resource.resource_id, exc)
            continue
        state[resource.resource_id] = {
            "dataset_id": resource.dataset_id,
            "dataset_title": resource.dataset_title,
            "resource_name": resource.resource_name,
            "resource_url": resource.resource_url,
            "summary_file": summary_filename,
            "processed_at": datetime.now(timezone.utc).isoformat(),
        }
        state_changed = True
        processed_count += 1
        if args.limit and processed_count >= args.limit:
            logging.info("Reached processing limit of %s resources", args.limit)
            break

    if state_changed:
        save_state(state_path, state)
        logging.info("Processed %s new resources", processed_count)
    else:
        logging.info("No new resources to process")

    return 0


if __name__ == "__main__":
    sys.exit(main())
